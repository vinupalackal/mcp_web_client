"""Generate MCP Client Web — product presentation (16 slides)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Palette ────────────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT  = RGBColor(0x3B, 0x82, 0xF6)
ACCENT2 = RGBColor(0x10, 0xB9, 0x81)
WARN    = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
PINK    = RGBColor(0xEC, 0x48, 0x99)
CYAN    = RGBColor(0x06, 0xB6, 0xD4)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xE2, 0xE8, 0xF0)
MID     = RGBColor(0x94, 0xA3, 0xB8)
BOX_BG  = RGBColor(0x1E, 0x29, 0x3B)
DIVIDER = RGBColor(0x33, 0x41, 0x55)

W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=DARK_BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def tx(slide, text, l, t, w, h, size=Pt(13), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def body_tf(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def card(slide, title, bullets, l, t, w, h,
         tc=ACCENT, bc=LIGHT, bs=Pt(11), bg_c=BOX_BG, bdr=DIVIDER):
    rect(slide, l, t, w, h, fill=bg_c, line=bdr, lw=Pt(0.75))
    tx(slide, title,
       l + Inches(0.12), t + Inches(0.1),
       w - Inches(0.24), Inches(0.35),
       size=Pt(12.5), bold=True, color=tc)
    tf = body_tf(slide,
                 l + Inches(0.14), t + Inches(0.48),
                 w - Inches(0.28), h - Inches(0.55))
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = "- " + b
        r.font.size = bs
        r.font.color.rgb = bc


def heading(slide, text, top=Inches(0.28), size=Pt(30), color=WHITE):
    tx(slide, text, Inches(0.6), top, Inches(12), Inches(0.55),
       size=size, bold=True, color=color)


def divider(slide, top):
    r = slide.shapes.add_shape(1, Inches(0.6), top, Inches(12.13), Pt(1))
    r.fill.solid()
    r.fill.fore_color.rgb = DIVIDER
    r.line.fill.background()


# ── Build ──────────────────────────────────────────────────────────────────
prs = new_prs()

# ─── SLIDE 1: TITLE ────────────────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(0.08), fill=ACCENT)
rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=ACCENT2)
tx(s, "MCP Client Web",
   Inches(1.0), Inches(1.7), Inches(11.3), Inches(1.3),
   size=Pt(54), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, "Browser-Based AI-Powered Device Diagnostics via Model Context Protocol",
   Inches(1.0), Inches(3.1), Inches(11.3), Inches(0.6),
   size=Pt(19), color=LIGHT, align=PP_ALIGN.CENTER)
rect(s, Inches(4.4), Inches(3.82), Inches(4.55), Pt(2), fill=ACCENT)
tx(s, "v0.4.0  |  FastAPI + Vanilla JS  |  JSON-RPC 2.0",
   Inches(1.0), Inches(4.05), Inches(11.3), Inches(0.4),
   size=Pt(13), color=MID, align=PP_ALIGN.CENTER)
pills = [
    ("Multi-LLM", ACCENT),
    ("265+ Tools", ACCENT2),
    ("SSO Auth", PURPLE),
    ("Parallel Exec", WARN),
]
for i, (label, color) in enumerate(pills):
    lx = Inches(1.7 + i * 2.5)
    rect(s, lx, Inches(5.0), Inches(2.2), Inches(0.44), fill=color)
    tx(s, label, lx, Inches(5.0), Inches(2.2), Inches(0.44),
       size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, "March 2026",
   Inches(0.6), Inches(6.85), Inches(4), Inches(0.35), size=Pt(11), color=MID)
tx(s, "CONFIDENTIAL",
   Inches(9.0), Inches(6.85), Inches(4), Inches(0.35),
   size=Pt(11), color=MID, align=PP_ALIGN.RIGHT)


# ─── SLIDE 2: AGENDA ───────────────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "Agenda")
divider(s, Inches(1.15))
agenda = [
    ("01", "Product Overview & Goals"),
    ("02", "System Architecture"),
    ("03", "Core Modules"),
    ("04", "MCP Protocol & Tool Execution"),
    ("05", "LLM Integration & Providers"),
    ("06", "Query Intelligence & Routing"),
    ("07", "Parallel Execution & Split-Phase"),
    ("08", "Repeated Tool Execution"),
    ("09", "Enterprise LLM Gateway"),
    ("10", "SSO & Multi-User"),
    ("11", "UI Features & Tool Output Display"),
    ("12", "Security, Deployment & Use Cases"),
]
for ci, col_items in enumerate([agenda[:6], agenda[6:]]):
    lx = Inches(0.6 + ci * 6.35)
    for ri, (num, title) in enumerate(col_items):
        ty = Inches(1.35 + ri * 0.85)
        rect(s, lx, ty, Inches(0.5), Inches(0.5), fill=ACCENT)
        tx(s, num, lx, ty, Inches(0.5), Inches(0.5),
           size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tx(s, title, lx + Inches(0.6), ty + Inches(0.06),
           Inches(5.4), Inches(0.42), size=Pt(14), color=LIGHT)


# ─── SLIDE 3: PRODUCT OVERVIEW ─────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "Product Overview")
divider(s, Inches(1.15))
tx(s, ("MCP Client Web is a browser-based chat interface that connects to MCP servers, "
       "enabling AI-assisted device diagnostics, system monitoring, and tool-driven workflows "
       "through natural language -- with no build step required."),
   Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.75), size=Pt(13), color=LIGHT)
goals = [
    ("No Build Step",   "Pure Vanilla JS SPA -- no webpack or npm build"),
    ("Distributed",     "MCP servers, LLM, and client on separate machines"),
    ("Multi-Server",    "Connect to multiple MCP servers simultaneously"),
    ("Multi-LLM",       "OpenAI, Ollama, Enterprise (Comcast) and Mock"),
    ("SSO Ready",       "Azure AD & Google OIDC with per-user data isolation"),
    ("Parallel Exec",   "Concurrent tool dispatch via asyncio.gather()"),
]
cols_g = [ACCENT, ACCENT2, WARN, PURPLE, PINK, CYAN]
for i, (title, desc) in enumerate(goals):
    col = i % 3
    row = i // 3
    lx = Inches(0.6 + col * 4.2)
    ty = Inches(2.15 + row * 1.6)
    rect(s, lx, ty, Inches(3.9), Inches(1.42), fill=BOX_BG, line=cols_g[i], lw=Pt(0.75))
    tx(s, title, lx + Inches(0.12), ty + Inches(0.1),
       Inches(3.65), Inches(0.38), size=Pt(13), bold=True, color=cols_g[i])
    tx(s, desc, lx + Inches(0.12), ty + Inches(0.52),
       Inches(3.65), Inches(0.8), size=Pt(11), color=LIGHT)


# ─── SLIDE 4: SYSTEM ARCHITECTURE ──────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "System Architecture")
divider(s, Inches(1.15))
tiers = [
    ("Browser (Frontend)", ACCENT, [
        "app.js -- Chat UI, tool panels, synthesis",
        "settings.js -- Server / LLM / Tool config",
        "localStorage -- Config persistence & cache",
        "401 guard  |  SSO user menu  |  Themes",
    ]),
    ("FastAPI Backend", ACCENT2, [
        "REST API endpoints (/api/*)",
        "Query Classifier & Mode Routing",
        "MCP Manager (JSON-RPC 2.0 client)",
        "LLM Client Factory (adapters)",
        "Session Manager (in-memory)",
        "Auth Middleware + SQLite DB",
    ]),
    ("External Services", WARN, [
        "MCP Servers -- Tool Registry & Executor",
        "OpenAI / Ollama / Enterprise LLM",
        "Azure AD / Google (OIDC IdP)",
        "LAN / WAN multi-machine support",
    ]),
]
for i, (title, color, bullets) in enumerate(tiers):
    lx = Inches(0.5 + i * 4.23)
    rect(s, lx, Inches(1.28), Inches(3.98), Inches(5.8), fill=BOX_BG, line=color, lw=Pt(1.5))
    tx(s, title, lx + Inches(0.12), Inches(1.38), Inches(3.74), Inches(0.38),
       size=Pt(14), bold=True, color=color)
    rect(s, lx, Inches(1.78), Inches(3.98), Pt(1), fill=color)
    tf = body_tf(s, lx + Inches(0.18), Inches(1.92), Inches(3.62), Inches(5.0))
    for j, b in enumerate(bullets):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = "- " + b
        r.font.size = Pt(12)
        r.font.color.rgb = LIGHT
    if i < 2:
        tx(s, "<->",
           Inches(4.39 + i * 4.23), Inches(3.95),
           Inches(0.32), Inches(0.4),
           size=Pt(16), bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)


# ─── SLIDE 5: CORE MODULES ─────────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "Core Modules")
divider(s, Inches(1.15))
modules = [
    ("backend/main.py", ACCENT, [
        "FastAPI app & REST endpoints",
        "Query classifier & mode scoring",
        "Multi-turn LLM loop (max 8 turns)",
        "Split-phase tool dispatch",
        "Domain-aware tool narrowing",
        "Synthesis prompt orchestration",
    ]),
    ("backend/mcp_manager.py", ACCENT2, [
        "JSON-RPC 2.0 initialize handshake",
        "tools/list -- tool discovery",
        "tools/call -- tool execution",
        "Chunked catalog (get_tools_for_llm)",
        "Parallel & repeated exec support",
        "Server health monitoring",
    ]),
    ("backend/llm_client.py", WARN, [
        "LLMClientFactory",
        "OpenAILLMClient",
        "OllamaLLMClient",
        "EnterpriseLLMClient",
        "Tool result format per provider",
        "Message adapter layer",
    ]),
    ("backend/session_manager.py", PURPLE, [
        "In-memory sessions",
        "Message history per session",
        "Tool trace events",
        "Conversation summary builder",
        "User-scoped keying (SSO mode)",
    ]),
    ("backend/auth/", PINK, [
        "OIDC provider ABC",
        "Azure AD & Google providers",
        "PKCE code verifier helpers",
        "JWKS cache (offline JWT verify)",
        "JWT issue / validate / revoke",
    ]),
    ("backend/static/", CYAN, [
        "app.js (Chat UI v21)",
        "settings.js (Tabbed modal)",
        "style.css (3 themes)",
        "login.html (SSO page)",
        "tool-tester.html (debug page)",
    ]),
]
for i, (name, color, bullets) in enumerate(modules):
    col = i % 3
    row = i // 3
    card(s, name, bullets,
         Inches(0.5 + col * 4.28),
         Inches(1.28 + row * 2.92),
         Inches(4.0), Inches(2.72), tc=color, bs=Pt(10.5))


# ─── SLIDE 6: MCP PROTOCOL ─────────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "MCP Protocol & Tool Execution")
divider(s, Inches(1.15))
steps = [
    ("1 User\nMessage", ACCENT),
    ("2 LLM +\nTools", ACCENT2),
    ("3 Parse\nCalls", WARN),
    ("4 Parallel\nExec", PINK),
    ("5 Results\nto LLM", PURPLE),
    ("6 Synthesis\nAnswer", ACCENT2),
]
bw = Inches(1.65)
bh = Inches(1.1)
gap = Inches(0.22)
for i, (label, color) in enumerate(steps):
    lx = Inches(0.45) + i * (bw + gap)
    rect(s, lx, Inches(1.35), bw, bh, fill=color)
    tx(s, label, lx, Inches(1.35), bw, bh,
       size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        tx(s, "->", lx + bw, Inches(1.7), gap, Inches(0.4),
           size=Pt(14), bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
tx(s, "JSON-RPC 2.0 Initialize",
   Inches(0.5), Inches(2.65), Inches(5.6), Inches(0.35),
   size=Pt(13), bold=True, color=ACCENT)
jrpc = (
    '{\n'
    '  "jsonrpc": "2.0", "id": 1,\n'
    '  "method": "initialize",\n'
    '  "params": {\n'
    '    "protocolVersion": "2024-11-05",\n'
    '    "clientInfo": {\n'
    '      "name": "mcp-client-web",\n'
    '      "version": "1.0"\n'
    '  }}\n'
    '}'
)
rect(s, Inches(0.5), Inches(3.05), Inches(5.6), Inches(3.0),
     fill=RGBColor(0x0D, 0x13, 0x1F))
tx(s, jrpc, Inches(0.65), Inches(3.15), Inches(5.3), Inches(2.8),
   size=Pt(9.5), color=ACCENT2)
tx(s, "Key Behaviours",
   Inches(6.4), Inches(2.65), Inches(6.5), Inches(0.35),
   size=Pt(13), bold=True, color=ACCENT)
rect(s, Inches(6.4), Inches(3.05), Inches(6.5), Inches(3.0),
     fill=BOX_BG, line=DIVIDER, lw=Pt(0.75))
tf = body_tf(s, Inches(6.6), Inches(3.15), Inches(6.1), Inches(2.8))
facts = [
    "Tool IDs namespaced: server_alias__tool_name",
    "Max 8 tool calls per turn (MCP_MAX_TOOL_CALLS_PER_TURN)",
    "Results >12K chars truncated before LLM",
    "asyncio.gather() fires all calls concurrently",
    "Waits for slowest device (max of all latencies)",
    "Tool results labeled [tool_name] (Ollama compat)",
    "Deduplication: same tool+args called once only",
]
for j, f in enumerate(facts):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.space_before = Pt(3)
    r = p.add_run()
    r.text = "- " + f
    r.font.size = Pt(11)
    r.font.color.rgb = LIGHT
tx(s, "Tool naming: home_mcp_server__get_dmesg",
   Inches(6.4), Inches(6.35), Inches(6.5), Inches(0.4),
   size=Pt(12), bold=True, color=ACCENT2)


# ─── SLIDE 7: LLM PROVIDERS ────────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "LLM Integration & Providers")
divider(s, Inches(1.15))
providers = [
    ("OpenAI", ACCENT, [
        "Configurable base URL",
        "API key authentication",
        "Uses tool_call_id in responses",
        "Supports GPT-4o, GPT-4-turbo, etc.",
        "Standard /v1/chat/completions",
    ]),
    ("Ollama", ACCENT2, [
        "Local or remote LAN instance",
        "No API key required",
        "Tool name only (no tool_call_id)",
        "Results labeled [tool_name] prefix",
        "llama3, mistral, qwen2.5, etc.",
    ]),
    ("Enterprise (Comcast Gateway)", WARN, [
        "OAuth 2.0 bearer token auth",
        "Token from /v2/oauth/token",
        "In-memory token cache",
        "OpenAI-compat /chat/completions",
        "Model catalog from UI",
    ]),
    ("Mock", MID, [
        "Testing / offline CI mode",
        "Deterministic responses",
        "No network calls made",
        "Full test suite coverage",
    ]),
]
for i, (name, color, bullets) in enumerate(providers):
    col = i % 2
    row = i // 2
    card(s, name, bullets,
         Inches(0.5 + col * 6.45),
         Inches(1.28 + row * 2.88),
         Inches(6.18), Inches(2.68), tc=color, bs=Pt(12))


# ─── SLIDE 8: QUERY ROUTING ────────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "Query Intelligence & Routing")
divider(s, Inches(1.15))
modes = [
    ("direct_fact",     ACCENT,  "+6 direct route\n+2 single domain\nEg: uptime?"),
    ("targeted_status", ACCENT2, "+3 keyword match\n+2 domains\nEg: show kernel logs"),
    ("full_diagnostic", WARN,    "+4 diagnostic words\nEg: why is it slow?"),
    ("follow_up",       PURPLE,  "+4 follow-up pattern\nEg: what about now?"),
]
for i, (mode, color, desc) in enumerate(modes):
    lx = Inches(0.5 + i * 3.2)
    rect(s, lx, Inches(1.28), Inches(3.0), Inches(1.72), fill=color)
    tx(s, mode, lx, Inches(1.38), Inches(3.0), Inches(0.4),
       size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, desc, lx + Inches(0.12), Inches(1.82), Inches(2.76), Inches(1.0),
       size=Pt(10.5), color=WHITE, align=PP_ALIGN.CENTER)
tx(s, "Direct Query Routes -- bypass scoring, go straight to named tool",
   Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.38),
   size=Pt(13), bold=True, color=ACCENT)
routes = [
    ("free_memory",  "how much free memory / available memory"),
    ("uptime",       "uptime / how long running / since reboot"),
    ("cpu_usage",    "cpu usage / load average"),
    ("disk_usage",   "disk space / free disk"),
    ("wan_ip",       "wan ip / public ip / external ip"),
    ("kernel_logs",  "dmesg / kernel log / last N kernel lines / syslog"),
]
rect(s, Inches(0.5), Inches(3.58), Inches(12.3), Inches(2.3),
     fill=BOX_BG, line=DIVIDER, lw=Pt(0.75))
tf = body_tf(s, Inches(0.7), Inches(3.68), Inches(12.0), Inches(2.1))
for j, (rt, pat) in enumerate(routes):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.space_before = Pt(2)
    r = p.add_run()
    r.text = "- {:<14}  ->  {}".format(rt, pat)
    r.font.size = Pt(11)
    r.font.color.rgb = LIGHT
tx(s, "Domain-Aware Narrowing: maps logs/memory/cpu/network/disk/wifi/uptime "
      "to tool keywords -> filters 265 tools to relevant subset",
   Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9),
   size=Pt(11), color=LIGHT)


# ─── SLIDE 9: PARALLEL & SPLIT-PHASE ───────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "Parallel Execution & Split-Phase Dispatch")
divider(s, Inches(1.15))
tx(s, "Serial (before):",
   Inches(0.6), Inches(1.32), Inches(3), Inches(0.35), size=Pt(12), bold=True, color=WARN)
serial = [("tool1", 2.1), ("tool2", 1.55), ("tool3", 0.95)]
offsets = [0.0, 2.18, 2.18 + 1.63]
for i, (lbl, w) in enumerate(serial):
    lx = Inches(0.6 + offsets[i])
    rect(s, lx, Inches(1.72), Inches(w), Inches(0.42), fill=WARN)
    tx(s, lbl, lx, Inches(1.72), Inches(w), Inches(0.42),
       size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, "Total latency = t1 + t2 + t3",
   Inches(0.6), Inches(2.22), Inches(6), Inches(0.3), size=Pt(10), color=MID)
tx(s, "Parallel -- asyncio.gather() (after):",
   Inches(0.6), Inches(2.65), Inches(5), Inches(0.35), size=Pt(12), bold=True, color=ACCENT2)
for i, (lbl, w) in enumerate(serial):
    rect(s, Inches(0.6), Inches(3.08 + i * 0.54), Inches(w), Inches(0.42), fill=ACCENT2)
    tx(s, lbl, Inches(0.6), Inches(3.08 + i * 0.54), Inches(w), Inches(0.42),
       size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, "Total latency = max(t1, t2, t3)",
   Inches(0.6), Inches(4.78), Inches(6), Inches(0.3), size=Pt(10), color=MID)
card(s, "Split-Phase Dispatch", [
    "When tool catalog > tools_split_limit (default 128)",
    "LLM queried once per chunk (read-only conversation snapshot)",
    "All tool_calls across chunks merged and deduped",
    "Injected into Turn 0 -- single parallel MCP execution phase",
    "Early-stop for direct_fact at confidence >= 0.75",
    "Sequential or concurrent chunks (configurable per session)",
], Inches(6.5), Inches(1.28), Inches(6.4), Inches(2.95), tc=ACCENT)
card(s, "Synthesis Prompt", [
    "After tools execute, system prompt switches to synthesis mode",
    "Per-tool status table: tool_name [ok/err]: 120-char preview",
    "Instruction: do NOT call more tools, read results",
    "kB -> MB/GB translation directive",
    "Green UI panel when LLM reports no issues",
    "Checkmark Analysis label vs notepad Analysis label",
], Inches(6.5), Inches(4.35), Inches(6.4), Inches(2.8), tc=ACCENT2)


# ─── SLIDE 10: REPEATED EXEC ───────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "Repeated Tool Execution -- mcp_repeated_exec")
divider(s, Inches(1.15))
tx(s, "Virtual client-side tool injected into LLM catalog -- never sent to an MCP server. "
      "LLM calls it like any other tool.",
   Inches(0.6), Inches(1.22), Inches(12), Inches(0.38),
   size=Pt(12.5), color=MID, italic=True)
card(s, "Parameters", [
    "target_tool    -- namespaced MCP tool to repeat",
    "repeat_count   -- number of runs (1-10, mandatory)",
    "interval_ms    -- delay between runs >=0ms (mandatory)",
    "tool_arguments -- args passed to target each run",
], Inches(0.5), Inches(1.72), Inches(4.35), Inches(2.52), tc=ACCENT)
card(s, "Execution Flow", [
    "1. Validate repeat_count & interval_ms (error if missing)",
    "2. Verify target_tool exists in registry",
    "3. Run target_tool x N sequentially with sleep(interval_ms)",
    "4. Save each run output to file (MCP_OUTPUT_DIR)",
    "5. Build cross-run synthesis prompt",
    "6. LLM analyses trend across all N runs",
], Inches(4.98), Inches(1.72), Inches(4.35), Inches(2.52), tc=ACCENT2)
card(s, "Use Cases", [
    "Memory leak detection",
    "CPU spin monitoring over time",
    "File descriptor growth tracking",
    "Trend & regression analysis",
    "Multi-point data collection",
], Inches(9.46), Inches(1.72), Inches(3.45), Inches(2.52), tc=WARN)
tx(s, "Tool Schema (excerpt)",
   Inches(0.5), Inches(4.4), Inches(12), Inches(0.35),
   size=Pt(13), bold=True, color=ACCENT)
schema = (
    '{ "name": "mcp_repeated_exec",\n'
    '  "description": "Execute an MCP tool N times at a fixed interval.",\n'
    '  "required": ["target_tool", "repeat_count", "interval_ms"],\n'
    '  "properties": {\n'
    '    "repeat_count": { "type":"integer", "minimum":1, "maximum":10 },\n'
    '    "interval_ms":  { "type":"integer", "minimum":0 } } }'
)
rect(s, Inches(0.5), Inches(4.82), Inches(12.3), Inches(2.3),
     fill=RGBColor(0x0D, 0x13, 0x1F))
tx(s, schema, Inches(0.65), Inches(4.92), Inches(12.0), Inches(2.1),
   size=Pt(10), color=ACCENT2)


# ─── SLIDE 11: ENTERPRISE GATEWAY ──────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "Enterprise LLM Gateway (Comcast Model Gateway)")
divider(s, Inches(1.15))
card(s, "Architecture", [
    "Additive, non-breaking -- standard providers unchanged",
    "EnterpriseLLMClient speaks OpenAI /chat/completions",
    "POST /api/enterprise/token -- fetch OAuth 2.0 bearer token",
    "In-memory token cache (per-process, clears on restart)",
    "POST /api/enterprise/token/status -- inspect cached token",
    "POST /api/enterprise/models -- fetch model catalog",
], Inches(0.5), Inches(1.28), Inches(5.9), Inches(3.0), tc=ACCENT)
card(s, "OAuth 2.0 Token Flow", [
    "1. UI sends client_id + client_secret to /api/enterprise/token",
    "2. Backend POSTs to Enterprise /v2/oauth/token",
    "3. Access token cached in-memory on backend",
    "4. Every LLM request: Authorization: Bearer <token>",
    "5. Token masked in all logs -- never returned to UI",
], Inches(6.55), Inches(1.28), Inches(6.35), Inches(3.0), tc=WARN)
card(s, "Settings Modal UI", [
    "Mode toggle: Standard vs Enterprise",
    "Enterprise panel: gateway URL + model name",
    "Client credentials: password-type inputs",
    "Fetch Token button -- validates connectivity",
    "Token status badge: valid / expired / not-fetched",
], Inches(0.5), Inches(4.42), Inches(5.9), Inches(2.65), tc=ACCENT2)
card(s, "Security Guarantees", [
    "Credentials never logged or returned to frontend",
    "Bearer token never exposed beyond backend memory",
    "HTTPS enforced for all enterprise endpoints",
    "Token expiry tracked; badge warns before expiry",
], Inches(6.55), Inches(4.42), Inches(6.35), Inches(2.65), tc=PURPLE)


# ─── SLIDE 12: SSO & MULTI-USER ────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "SSO Authentication & Multi-User Support")
divider(s, Inches(1.15))
card(s, "Auth Layer", [
    "OIDC Authorization Code Flow with PKCE",
    "Providers: Azure AD & Google",
    "HttpOnly JWT session cookie (app_token)",
    "Middleware guard on every /api/* route",
    "JWKS cache for offline JWT verification",
    "Stateless -- no server-side session store",
], Inches(0.5), Inches(1.28), Inches(4.12), Inches(3.15), tc=ACCENT)
card(s, "Per-User Data Isolation", [
    "UserScopedServerStore -- per-user MCP servers",
    "UserScopedLLMConfigStore -- per-user LLM config",
    "Sessions keyed by (user_id, session_id)",
    "UserSettingsStore -- theme & preferences",
    "SQLite default, Postgres-ready via SQLAlchemy",
    "localStorage becomes read-ahead cache only",
], Inches(4.77), Inches(1.28), Inches(4.12), Inches(3.15), tc=ACCENT2)
card(s, "Admin Controls", [
    "GET /api/admin/users -- list all users",
    "PATCH /api/admin/users/{id} -- enable/disable",
    "SSO_ADMIN_EMAILS env var -- admin allowlist",
    "Role-based access on UserProfile",
], Inches(9.04), Inches(1.28), Inches(3.9), Inches(3.15), tc=WARN)
tx(s, "New API Endpoints",
   Inches(0.5), Inches(4.58), Inches(12), Inches(0.35),
   size=Pt(13), bold=True, color=ACCENT)
eps = [
    "/auth/login  |  /auth/callback/{provider}  |  /auth/logout",
    "/api/users/me  |  /api/users/me/settings",
    "/api/admin/users  |  /api/admin/users/{user_id}",
]
rect(s, Inches(0.5), Inches(4.98), Inches(12.3), Inches(1.7),
     fill=BOX_BG, line=DIVIDER, lw=Pt(0.75))
tf = body_tf(s, Inches(0.7), Inches(5.08), Inches(12.0), Inches(1.5))
for j, e in enumerate(eps):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.space_before = Pt(4)
    r = p.add_run()
    r.text = "- " + e
    r.font.size = Pt(12)
    r.font.color.rgb = ACCENT2


# ─── SLIDE 13: UI FEATURES ─────────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "UI Features & Tool Output Display")
divider(s, Inches(1.15))
card(s, "Chat Interface (app.js v21)", [
    "LibreChat-inspired layout -- header / chat / input",
    "Tool executions shown BEFORE synthesis answer",
    "Failed tools auto-expanded (no clicking needed)",
    "Header: Tool Results -- N executed, N failed",
    "Cleaned result_text (unwrapped MCP JSON-RPC envelope)",
    "Analysis label on synthesis block (green = no issues)",
    "Green panel when LLM reports no issues found",
], Inches(0.5), Inches(1.28), Inches(4.12), Inches(5.5), tc=ACCENT, bs=Pt(11))
card(s, "Settings Modal (settings.js)", [
    "Tabs: My Account | MCP Servers | LLM Config | Tools",
    "Server CRUD with inline validation & error messages",
    "Refresh Tools -> JSON-RPC tool discovery",
    "LLM toggle: Standard / Enterprise gateway",
    "Connection test per provider",
    "Tools tab: grouped by server alias",
    "Modal stays open during multi-step operations",
], Inches(4.77), Inches(1.28), Inches(4.12), Inches(5.5), tc=ACCENT2, bs=Pt(11))
card(s, "Themes & Debug Tools", [
    "Three themes: Light / Dark / Deep Teal",
    "Theme saved per-user (SSO) or localStorage",
    "pagehide/visibilitychange -> save session state",
    "Chat view state restored on tab reload",
    "tool-tester.html -- test MCP tools in isolation",
    "Emoji console logs for quick filtering",
    "Swagger UI at /docs, ReDoc at /redoc",
], Inches(9.04), Inches(1.28), Inches(3.9), Inches(5.5), tc=WARN, bs=Pt(11))


# ─── SLIDE 14: SECURITY & DEPLOYMENT ───────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "Security & Deployment")
divider(s, Inches(1.15))
card(s, "Security Constraints", [
    "HTTPS enforced for all MCP server URLs",
    "MCP_ALLOW_HTTP_INSECURE=true for local dev only",
    "API keys / tokens never logged (masked)",
    "Password-type inputs for all credential fields",
    "CORS restricted to known origins in production",
    "JWT PKCE mandatory -- no implicit OAuth flow",
    "Admin allowlist via SSO_ADMIN_EMAILS env var",
], Inches(0.5), Inches(1.28), Inches(4.0), Inches(5.6), tc=WARN)
card(s, "Key Environment Variables", [
    "MCP_ALLOW_HTTP_INSECURE   (false)",
    "MCP_REQUEST_TIMEOUT_MS    (20000)",
    "MCP_MAX_TOOL_CALLS_PER_TURN (8)",
    "MCP_MAX_TOOLS_PER_REQUEST   (128)",
    "MCP_MAX_TOOL_OUTPUT_CHARS_TO_LLM (12000)",
    "MCP_SPLIT_PHASE_MIN_CONFIDENCE (0.75)",
    "OPENAI_API_KEY / OPENAI_BASE_URL",
    "OLLAMA_BASE_URL",
    "SECRET_KEY  (JWT signing for SSO)",
], Inches(4.65), Inches(1.28), Inches(4.3), Inches(5.6), tc=ACCENT, bs=Pt(11))
card(s, "Deployment Options", [
    "Single machine: all on localhost:8000",
    "Multi-machine: MCP on LAN IP, LLM on separate host",
    "Enterprise SSO: SQLite or Postgres, OIDC redirect URIs",
    "Python 3.9+ compatible",
    "No build step -- static files served direct",
    "uvicorn backend.main:app --host 0.0.0.0",
    "OpenAPI docs: /docs  |  /redoc",
], Inches(9.1), Inches(1.28), Inches(3.85), Inches(5.6), tc=ACCENT2, bs=Pt(11))


# ─── SLIDE 15: USE CASES ───────────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(1.1), fill=BOX_BG)
heading(s, "Use Cases")
divider(s, Inches(1.15))
cases = [
    ("Device Diagnostics",
     "Natural language queries like 'how much free memory?' or 'show disk usage' "
     "route directly to the right tool. Results translated to human-readable form.",
     ACCENT),
    ("Kernel Log Analysis",
     "'Show me the last 100 kernel log lines for errors.' "
     "Routes to dmesg tool, auto-expands failed results. Green panel = all clear.",
     ACCENT2),
    ("Longitudinal Monitoring",
     "'Sample memory every 30s, 5 times and check for leaks.' "
     "mcp_repeated_exec runs N times sequentially, LLM synthesizes trend.",
     WARN),
    ("Root Cause Analysis",
     "'Why is this device slow?' Full diagnostic mode: broad tool selection, "
     "classification pass, then targeted multi-tool execution with synthesis.",
     PURPLE),
    ("Enterprise Multi-User",
     "SSO-authenticated deployment. Each engineer has isolated MCP & LLM configs. "
     "Comcast Model Gateway for on-prem LLM inference via OAuth 2.0.",
     PINK),
    ("Tool Development & Testing",
     "Use tool-tester.html to validate individual MCP tools in isolation. "
     "Swagger /docs for API exploration without the full chat UI.",
     CYAN),
]
for i, (title, desc, color) in enumerate(cases):
    col = i % 2
    row = i // 2
    lx = Inches(0.5 + col * 6.52)
    ty = Inches(1.28 + row * 1.98)
    rect(s, lx, ty, Inches(6.25), Inches(1.82), fill=BOX_BG, line=color, lw=Pt(1.5))
    rect(s, lx, ty, Inches(0.22), Inches(1.82), fill=color)
    tx(s, title, lx + Inches(0.32), ty + Inches(0.12),
       Inches(5.8), Inches(0.38), size=Pt(13), bold=True, color=color)
    tx(s, desc, lx + Inches(0.32), ty + Inches(0.55),
       Inches(5.8), Inches(1.15), size=Pt(11), color=LIGHT)


# ─── SLIDE 16: SUMMARY ─────────────────────────────────────────────────────
s = add_slide(prs)
bg(s)
rect(s, 0, 0, W, Inches(0.08), fill=ACCENT)
rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=ACCENT2)
tx(s, "Summary",
   Inches(1.0), Inches(0.78), Inches(11.3), Inches(0.65),
   size=Pt(38), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
summary_items = [
    (ACCENT,  "Multi-LLM",          "OpenAI  |  Ollama  |  Enterprise (Comcast)  |  Mock"),
    (ACCENT2, "Tool Execution",     "JSON-RPC 2.0  |  Parallel  |  Repeated  |  Split-Phase"),
    (WARN,    "Query Intelligence", "4 routing modes  |  Direct routes  |  Domain narrowing"),
    (PURPLE,  "Auth & Multi-User",  "Azure AD  |  Google OIDC  |  PKCE  |  Per-user isolation"),
    (PINK,    "UI / UX",            "LibreChat-inspired  |  Dark/Teal/Light  |  Tool panels"),
    (CYAN,    "DevOps",             "No build step  |  Python 3.9+  |  OpenAPI  |  SQLite/PG"),
]
for i, (color, title, desc) in enumerate(summary_items):
    row = i % 3
    col = i // 3
    lx = Inches(0.6 + col * 6.55)
    ty = Inches(1.62 + row * 1.65)
    rect(s, lx, ty, Inches(0.1), Inches(1.3), fill=color)
    tx(s, title, lx + Inches(0.25), ty + Inches(0.05),
       Inches(6.0), Inches(0.4), size=Pt(15), bold=True, color=color)
    tx(s, desc, lx + Inches(0.25), ty + Inches(0.52),
       Inches(6.0), Inches(0.68), size=Pt(13), color=LIGHT)
tx(s, "github.com/vinupalackal/mcp_web_client",
   Inches(0.6), Inches(6.88), Inches(12.13), Inches(0.35),
   size=Pt(12), color=MID, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────
OUT = "/Users/vinupalackal/Project/MCP/mcp_client/MCP_Client_Web_Presentation.pptx"
prs.save(OUT)
print(f"Saved {len(prs.slides)} slides -> {OUT}")
